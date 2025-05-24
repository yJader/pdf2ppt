from pathlib import Path
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
import os

import typer

pdf2ppt_app = typer.Typer()


def extract_pdf_comments_with_pages(pdf_path: Path) -> dict[int, list[str]]:
    """
    从 PDF 文件中提取注释及其对应的页码。
    假定 pdfcomment 生成的是标准的文本注释。

    参数:
        pdf_path (Path): PDF 文件的路径。

    返回:
        dict: 一个字典，键是页码 (0-indexed)，值是该页上注释内容的列表。
              例如: {0: ["注释1内容", "注释2内容"], 1: ["另一页的注释"]}
    """
    comments_by_page = {}
    try:
        doc = fitz.open(pdf_path)
    except Exception as e:
        print(f"错误：无法打开 PDF 文件 '{pdf_path}'. {e}")
        return comments_by_page

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        page_comments = []
        # 遍历页面上的所有注释 (annotations)
        for annot in page.annots():
            # pdfcomment 通常创建 'Text' 类型的注释
            # annot.type[0] 是注释类型的数字代码, annot.type[1] 是名称
            # 'Text' (sticky note) 通常是类型 0 或 8 (根据 PDF 标准版本)
            # 我们主要关心 annot.info["content"]
            if (
                annot.type[1] == "Text" or annot.type[0] == 0 or annot.type[0] == 8
            ):  # 尝试兼容不同情况
                content = annot.info.get("content", "")
                if content:  # 确保有内容
                    # pdfcomment 可能在内容中加入作者等信息，这里简单提取主要内容
                    # 例如，如果注释格式是 "Author: \nComment text"，我们可能需要进一步处理
                    # 但通常 annot.info["content"] 直接就是注释文本
                    page_comments.append(content)
            # 你可以取消下面这行的注释来调试和查看所有注释的类型和信息
            # print(f"Page {page_num}, Annot Type: {annot.type}, Info: {annot.info}")

        if page_comments:
            comments_by_page[page_num] = page_comments

    doc.close()
    return comments_by_page


def convert_pdf_to_ppt_with_comments(
    pdf_path: Path,
    ppt_path: Path,
    comments_data: dict[int, list[str]],
    output_dpi: int = 600,
):
    """
    将 PDF 转换为 PPT，并将提取的注释添加到幻灯片备注中。

    参数:
        pdf_path (Path): 输入的 PDF 文件路径。
        ppt_path (Path): 输出的 PPT 文件路径。
        comments_data (dict): 从 extract_pdf_comments_with_pages 获取的注释数据。
        output_dpi (int): 输出图片的 DPI (dots per inch)，可以根据需要调整。
    """
    try:
        doc = fitz.open(pdf_path)
    except Exception as e:
        print(f"错误：无法打开 PDF 文件 '{pdf_path}' 进行转换. {e}")
        return

    prs = Presentation()

    # 获取 PDF 首页的尺寸以设置 PPT 幻灯片尺寸 (可选, 但推荐)
    # PyMuPDF 使用点 (points) 作为单位, 1 inch = 72 points
    # PowerPoint 使用 EMU (English Metric Units), python-pptx 用 Inches 辅助
    if len(doc) > 0:
        first_page = doc.load_page(0)
        pdf_width_pt = first_page.rect.width
        pdf_height_pt = first_page.rect.height

        # 设置 PPT 幻灯片尺寸以匹配 PDF 宽高比
        # 默认的 PPT 尺寸是 10x7.5 inches (4:3) 或 10x5.625 inches (16:9)
        # 你可以根据需要调整，或者直接使用默认然后让图片缩放
        prs.slide_width = Inches(pdf_width_pt / 72.0)
        prs.slide_height = Inches(pdf_height_pt / 72.0)

    # 选择一个空白的幻灯片布局 (通常索引为 5 或 6)
    # 你可以通过 `for i, layout in enumerate(prs.slide_layouts): print(i, layout.name)` 查看所有可用布局
    try:
        blank_slide_layout = prs.slide_layouts[6]  # 索引 6 通常是 "Blank"
    except IndexError:
        print("警告: 找不到索引为6的空白幻灯片布局，使用第一个布局。")
        blank_slide_layout = prs.slide_layouts[0]

    temp_image_dir = "temp_pdf_images"
    if not os.path.exists(temp_image_dir):
        os.makedirs(temp_image_dir)

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)

        # 1. 将 PDF 页面转换为图片
        # 可以调整 dpi (dots per inch) 来控制图片质量和大小
        pix = page.get_pixmap(dpi=output_dpi, annots=False)
        image_filename = os.path.join(temp_image_dir, f"page_{page_num + 1}.png")
        pix.save(image_filename)

        # 2. 在 PPT 中添加新幻灯片并将图片插入
        slide = prs.slides.add_slide(blank_slide_layout)

        # 添加图片并使其填充整个幻灯片
        # left, top, width, height
        # 如果幻灯片尺寸已设为 PDF 页面尺寸，则 left=0, top=0, width=prs.slide_width, height=prs.slide_height
        pic = slide.shapes.add_picture(
            image_filename,
            Inches(0),
            Inches(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )

        # 3. 添加注释到幻灯片备注
        if page_num in comments_data:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            # 清除可能存在的默认文本
            text_frame.clear()
            p = text_frame.paragraphs[0]  # 获取第一个段落
            p.text = "\n".join(comments_data[page_num])  # 将所有注释合并，用换行符分隔
            print(f"添加备注: {comments_data[page_num]} 到第 {page_num + 1} 页")

            # 如果需要将每个注释作为单独的段落：
            # text_frame.text = comments_data[page_num][0]
            # for i in range(1, len(comments_data[page_num])):
            #     p = text_frame.add_paragraph()
            #     p.text = comments_data[page_num][i]

        print(f"处理完成：第 {page_num + 1} 页 / 共 {len(doc)} 页")

    # 清理临时图片文件
    for page_num in range(len(doc)):
        image_filename = os.path.join(temp_image_dir, f"page_{page_num + 1}.png")
        if os.path.exists(image_filename):
            os.remove(image_filename)
    if os.path.exists(temp_image_dir) and not os.listdir(
        temp_image_dir
    ):  # 检查目录是否为空
        os.rmdir(temp_image_dir)
    elif os.path.exists(temp_image_dir):
        print(f"警告: 临时图片目录 {temp_image_dir} 未被完全清空。")

    doc.close()
    try:
        prs.save(ppt_path)
        print(f"成功！PPT 文件已保存到: {ppt_path}")
    except Exception as e:
        print(f"错误：无法保存 PPT 文件 '{ppt_path}'. {e}")


@pdf2ppt_app.command()
def convert(
    pdf_input_path: Path = typer.Option(
        ..., "-i", "--pdf-input-path", help="输入的 PDF 文件路径"
    ),
    ppt_output_path: Path = typer.Option(
        ..., "-o", "--ppt-output-path", help="输出的 PPT 文件路径"
    ),
    output_dpi: int = typer.Option(
        600, "-d", "--dpi", help="输出图片的 DPI (dots per inch)"
    ),
):
    """
    将 PDF 文件转换为 PPT 文件，并提取注释添加到备注中。

    参数:
        pdf_input_path (str): 输入的 PDF 文件路径。
        ppt_output_path (str): 输出的 PPT 文件路径。
    """
    if not pdf_input_path.exists():
        print(f"错误: 输入的 PDF 文件 '{pdf_input_path}' 不存在。请检查路径。")
        return

    if ppt_output_path is None:
        ppt_output_path = Path("output") / pdf_input_path.with_suffix(".pptx")
        print(f"输出的 PPT 文件路径未指定，使用默认路径: {ppt_output_path}")

    ppt_output_path.parent.mkdir(parents=True, exist_ok=True)

    print(f"正在从 '{pdf_input_path}' 提取注释...")
    extracted_comments = extract_pdf_comments_with_pages(pdf_input_path)

    if not extracted_comments:
        print("未在 PDF 中找到注释，或者提取失败。仍将尝试转换页面。")
    else:
        print(f"提取到 {sum(len(v) for v in extracted_comments.values())} 条注释。")

    print("正在将 PDF 转换为 PPT (带注释)...")
    convert_pdf_to_ppt_with_comments(
        pdf_input_path, ppt_output_path, extracted_comments, output_dpi
    )


if __name__ == "__main__":
    pdf2ppt_app()
